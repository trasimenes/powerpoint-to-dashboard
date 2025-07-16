import json
import os
import sqlite3
from datetime import datetime
from pathlib import Path
from tempfile import NamedTemporaryFile

from flask import Flask, flash, redirect, render_template, request
from pptx import Presentation

DB_PATH = Path("database.db")

app = Flask(__name__)
app.secret_key = "secret-key"  # In production, use env variable


def init_db():
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            """CREATE TABLE IF NOT EXISTS extractions (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp TEXT,
                filename TEXT,
                slide_start INTEGER,
                slide_end INTEGER,
                kpi TEXT,
                table_data TEXT
            )"""
        )


def insert_record(filename, slide_start, slide_end, kpi, table_data):
    with sqlite3.connect(DB_PATH) as conn:
        conn.execute(
            "INSERT INTO extractions (timestamp, filename, slide_start, slide_end, kpi, table_data) VALUES (?, ?, ?, ?, ?, ?)",
            (
                datetime.utcnow().isoformat(sep=" ", timespec="seconds"),
                filename,
                slide_start,
                slide_end,
                json.dumps(kpi, ensure_ascii=False),
                json.dumps(table_data, ensure_ascii=False),
            ),
        )
        conn.commit()


def get_history():
    with sqlite3.connect(DB_PATH) as conn:
        cursor = conn.execute(
            "SELECT timestamp, filename, slide_start, slide_end, kpi FROM extractions ORDER BY id DESC"
        )
        rows = cursor.fetchall()
    history = [
        {
            "timestamp": r[0],
            "filename": r[1],
            "slide_start": r[2],
            "slide_end": r[3],
            "kpi": json.loads(r[4]),
        }
        for r in rows
    ]
    return history


def parse_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                texts.append(text)
    return texts


def parse_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            headers = [cell.text.strip() for cell in table.rows[0].cells]
            rows = []
            for row in table.rows[1:]:
                rows.append([cell.text.strip() for cell in row.cells])
            return {"headers": headers, "rows": rows}
    return {"headers": [], "rows": []}


def extract_pptx(path, slide_start, slide_end):
    prs = Presentation(path)
    slides = prs.slides
    if slide_start - 1 >= len(slides) or slide_end - 1 >= len(slides):
        raise ValueError("Slide range out of bounds")
    kpi_slide = slides[slide_start - 1]
    table_slide = slides[slide_end - 1]
    kpis = parse_slide_text(kpi_slide)
    table = parse_table(table_slide)
    return kpis, table


@app.route("/", methods=["GET", "POST"])
def upload():
    if request.method == "POST":
        if "pptx" not in request.files:
            flash("No file part")
            return redirect(request.url)
        file = request.files["pptx"]
        if file.filename == "":
            flash("No selected file")
            return redirect(request.url)
        slide_start = int(request.form.get("start", 31))
        slide_end = int(request.form.get("end", 32))
        with NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name
        try:
            kpis, table = extract_pptx(tmp_path, slide_start, slide_end)
            insert_record(file.filename, slide_start, slide_end, kpis, table)
            slides = (slide_start, slide_end)
            return render_template(
                "dashboard.html", kpis=kpis, table=table, slides=slides
            )
        except Exception as e:
            flash(f"Error processing PPTX: {e}")
            return redirect(request.url)
        finally:
            os.unlink(tmp_path)
    return render_template("upload.html")


@app.route("/history")
def history():
    history = get_history()
    return render_template("history.html", history=history)


if __name__ == "__main__":
    init_db()
    app.run(debug=True)
