"""
cpfr_pptx_parser.py

Extraction "large" des données de la slide CPFR (ex: slide 31) à partir d'un .pptx.

Renvoie un dictionnaire structuré conforme à l'`ingest_payload_schema`
défini pour l'app CPFR Weekly Dashboard.

Usage direct (CLI):
    python cpfr_pptx_parser.py path/to/file.pptx --slide 31 --week-start 2025-07-14

Dans le code (import):
    from cpfr_pptx_parser import parse_cpfr_slide
    data = parse_cpfr_slide("deck.pptx", slide_number=31, week_start_date="2025-07-14")
"""

import re
import json
import argparse
from datetime import date
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple

from pptx import Presentation
from unidecode import unidecode

# -------------------------------
# Helpers: numeric parsing
# -------------------------------

_NUM_RE = re.compile(r"([+-]?\d+(?:[.,]\d+)?)([KkMm%€]?)")

def _normalize_number_fragment(raw: str) -> float:
    """Normalize a numeric fragment like '2,27M', '342K', '917', '0,53', '118k€' (case insensitive).

    Returns float (not divided if %, the caller handles).
    """
    txt = raw.strip()
    txt = txt.replace(" ", "")  # remove spaces in 2 475
    txt = txt.replace("\xa0", "")
    # separate suffix manually if double
    suffix = ""
    if txt.endswith(("€", "E", "e")):
        suffix = "€"
        txt = txt[:-1]
    # capture possible K/M
    mult = 1
    if txt.lower().endswith("k"):
        mult = 1_000
        txt = txt[:-1]
    elif txt.lower().endswith("m"):
        mult = 1_000_000
        txt = txt[:-1]
    txt = txt.replace(",", ".")
    try:
        val = float(txt) * mult
    except ValueError:
        val = 0.0
    return val

def parse_currency(raw: str) -> Optional[float]:
    m = _NUM_RE.search(raw.replace(" ", ""))
    if not m:
        return None
    return _normalize_number_fragment(m.group(1) + (m.group(2) if m.group(2) in ("K","k","M","m") else ""))

def parse_percent(raw: str) -> Optional[float]:
    """Return decimal (0.06 for +6%)."""
    txt = raw.strip()
    sign = 1
    if txt.startswith("+"): sign = 1; txt = txt[1:]
    elif txt.startswith("-"): sign = -1; txt = txt[1:]
    txt = txt.replace("%", "")
    txt = txt.replace(" ", "")
    txt = txt.replace(",", ".")
    try:
        return sign * (float(txt) / 100.0)
    except ValueError:
        return None

def parse_number_like(raw: str) -> Optional[float]:
    """Generic numeric (currency or plain). No percent division."""
    return parse_currency(raw)


# -------------------------------
# Text scanning utilities
# -------------------------------

def shape_text_iter(slide):
    """Yield (idx, text, norm_text) for each shape w/ text."""
    for i, sh in enumerate(slide.shapes):
        if not hasattr(sh, "text"):
            continue
        t = sh.text.strip()
        if not t:
            continue
        yield i, t, unidecode(t).lower()


def _find_shape(slide, contains_tokens: List[str]) -> Optional[str]:
    toks = [unidecode(tok).lower() for tok in contains_tokens]
    for _, t, norm in shape_text_iter(slide):
        if all(tok in norm for tok in toks):
            return t
    return None


# -------------------------------
# KPI Header Parsing
# -------------------------------

HEADER_PATTERNS = {
    "sessions": r"(\d[\d\s.,]*[KkMm]?)\s*(?:Nb\s*of\s*sessions|sessions)?",
    "revenue_b2c": r"(\d[\d\s.,]*[KkMm]?)\s*M?€",
    "average_basket_value": r"(\d[\d\s.,]*)(?:€|e)",
    "conversion_rate": r"(\d[\d\s.,]*)\s*%",
    "nb_bookings": r"(\d[\d\s.,]*)\s*(?:Nb\s*of\s*bookings|bookings)"
}

def parse_kpi_header(slide_text_all: str) -> Dict[str, Any]:
    """
    slide_text_all: big concatenated string of all header text shapes.
    Returns dict of KPIs; None where missing.
    """
    txt = slide_text_all
    data = {}

    # Sessions
    m = re.search(r"(\d[\d\s.,]*[KkMm]?)\s*Nb\s*of\s*sessions", txt, flags=re.I)
    data["sessions"] = _normalize_number_fragment(m.group(1)) if m else None

    # Revenue
    m = re.search(r"(\d[\d\s.,]*)(?:M)?\s*€\s*Web\s*B2C", txt, flags=re.I)
    if m:
        val = _normalize_number_fragment(m.group(1) + "M") if "M€" in txt[m.start():m.end()+2] else _normalize_number_fragment(m.group(1))
        # Heuristic: if number < 10k and "M€" around, multiply
        if val < 10_000 and "M€" in txt[m.start()-5:m.end()+5]:
            val *= 1_000_000
        data["revenue_b2c"] = val
    else:
        data["revenue_b2c"] = None

    # ABV
    m = re.search(r"Average\s*basket\s*value.*?(\d[\d\s.,]*)\s*€", txt, flags=re.I|re.S)
    data["average_basket_value"] = _normalize_number_fragment(m.group(1)) if m else None

    # CR
    m = re.search(r"Conversion\s*rate.*?(\d[\d\s.,]*)\s*%", txt, flags=re.I|re.S)
    data["conversion_rate"] = parse_percent(m.group(1)+"%") if m else None

    # Bookings
    m = re.search(r"(\d[\d\s.,]*)\s*Nb\s*of\s*bookings", txt, flags=re.I)
    if m:
        data["nb_bookings"] = int(_normalize_number_fragment(m.group(1)))
    else:
        data["nb_bookings"] = None

    # Variation blocks: capture +6% VS LY etc.
    # We'll pattern: ([+-]\s*\d+[.,]?\d*)%\s*VS\s*(LY|LW)
    vs = re.findall(r"([+-]\s*\d+[.,]?\d*)%\s*VS\s*(LY|LW)", txt, flags=re.I)
    # vs returns list in order; we need to map by KPI based on prox detection -> fallback: user populates later
    # For robustness we parse per KPI region:
    data.update(_parse_variations_per_kpi(txt))

    return data


def _parse_variations_per_kpi(txt: str) -> Dict[str, Any]:
    """
    More accurate: locate each KPI label region then parse +x% vs LY / LW in that substring window.
    """
    results = {
        "vs_ly_sessions": None, "vs_lw_sessions": None,
        "vs_ly_revenue": None,  "vs_lw_revenue": None,
        "vs_ly_abv": None,      "vs_lw_abv": None,
        "vs_ly_cr": None,       "vs_lw_cr": None,
        "vs_ly_bookings": None, "vs_lw_bookings": None,
    }

    def grab_window(label_regex: str) -> str:
        m = re.search(label_regex, txt, flags=re.I)
        if not m:
            return ""
        start = max(0, m.start()-40)
        end   = min(len(txt), m.end()+40)
        return txt[start:end]

    win_sessions = grab_window(r"Nb\s*of\s*sessions")
    win_revenue  = grab_window(r"Web\s*B2C\s*Global\s*revenue")
    win_abv      = grab_window(r"Average\s*basket\s*value")
    win_cr       = grab_window(r"Conversion\s*rate")
    win_book     = grab_window(r"Nb\s*of\s*bookings")

    def cap_variations(window: str):
        return re.findall(r"([+-]\s*\d+[.,]?\d*)%\s*VS\s*(LY|LW)", window, flags=re.I)

    for raw, tag in cap_variations(win_sessions):
        pc = parse_percent(raw+"%")
        if tag.upper()=="LY": results["vs_ly_sessions"]=pc
        else: results["vs_lw_sessions"]=pc

    for raw, tag in cap_variations(win_revenue):
        pc = parse_percent(raw+"%")
        if tag.upper()=="LY": results["vs_ly_revenue"]=pc
        else: results["vs_lw_revenue"]=pc

    for raw, tag in cap_variations(win_abv):
        pc = parse_percent(raw+"%")
        if tag.upper()=="LY": results["vs_ly_abv"]=pc
        else: results["vs_lw_abv"]=pc

    for raw, tag in cap_variations(win_cr):
        pc = parse_percent(raw+"%")
        if tag.upper()=="LY": results["vs_ly_cr"]=pc
        else: results["vs_lw_cr"]=pc

    for raw, tag in cap_variations(win_book):
        pc = parse_percent(raw+"%")
        if tag.upper()=="LY": results["vs_ly_bookings"]=pc
        else: results["vs_lw_bookings"]=pc

    return results


# -------------------------------
# Overview block: best traffic day
# -------------------------------

def parse_overview_block(txt: str) -> Dict[str, Any]:
    """
    Parse 'OVERVIEW PERFORMANCES' shape text.
    Grabs best day sessions & revenue.
    """
    data = {"best_day_sessions": None, "best_day_revenue": None, "raw_overview_text": txt}

    # best day line
    m = re.search(r"Best\s*traffic\s*/?\s*revenue\s*day.*?(\d{1,2}\w*\s*\w+).*?(\d[\d\s.,]*[KkMm]?)\s*sessions.*?(\d[\d\s.,]*[KkMm]?)\s*€", txt, flags=re.I)
    if m:
        # date text m.group(1) ignored (we don't store date here in schema; could extend)
        data["best_day_sessions"] = int(_normalize_number_fragment(m.group(2)))
        data["best_day_revenue"] = _normalize_number_fragment(m.group(3))
        return data

    # fallback: capture 2 numbers K inside same line
    m = re.search(r"Best.*?(\d[\d\s.,]*[KkMm]?).*?sessions.*?(\d[\d\s.,]*[KkMm]?)\s*€", txt, flags=re.I)
    if m:
        data["best_day_sessions"] = int(_normalize_number_fragment(m.group(1)))
        data["best_day_revenue"] = _normalize_number_fragment(m.group(2))
    return data


# -------------------------------
# Offers block
# -------------------------------

def parse_offers_block(txt: str) -> Dict[str, Any]:
    """
    Parse 'FOCUS OFFERS' bullet region.
    """
    data = {
        "last_minute_pct": None,
        "early_booking_pct": None,
        "summer_flash_revenue": None,
        "summer_flash_bookings": None,
        "summer_flash_abv": None,
        "lead_gen_revenue": None,
        "lead_gen_bookings": None,
        "raw_offers_text": txt
    }

    # Last Minute %
    m = re.search(r"(\d+[.,]?\d*)%\s*bookings\s*on\s*Last\s*Minute", txt, flags=re.I)
    if m: data["last_minute_pct"] = parse_percent(m.group(1)+"%")

    # Early Booking %
    m = re.search(r"(\d+[.,]?\d*)%\s*bookings\s*on\s*Early\s*Booking", txt, flags=re.I)
    if m: data["early_booking_pct"] = parse_percent(m.group(1)+"%")

    # Summer Flash Sale revenue
    # Example "Summer Flash Sale : 1,4M€ (60% of total revenue), 1,4K booking & 924€ ABV."
    m = re.search(r"Summer\s*Flash\s*Sale\s*:\s*([\d\s.,]*[KkMm]?)\s*€", txt, flags=re.I)
    if m: data["summer_flash_revenue"] = _normalize_number_fragment(m.group(1))

    # bookings
    m = re.search(r"Flash\s*Sale.*?([\d\s.,]*[KkMm]?)\s*book", txt, flags=re.I)
    if m: data["summer_flash_bookings"] = int(_normalize_number_fragment(m.group(1)))

    # ABV
    m = re.search(r"(\d[\d\s.,]*)\s*€\s*ABV", txt, flags=re.I)
    if m: data["summer_flash_abv"] = _normalize_number_fragment(m.group(1))

    # Lead gen revenue
    m = re.search(r"Lead\s*gen\s*:\s*([\d\s.,]*[KkMm]?)\s*€", txt, flags=re.I)
    if m: data["lead_gen_revenue"] = _normalize_number_fragment(m.group(1))

    # Lead gen bookings not explicit -> derive 1% of total bookings? not safe; we skip unless pattern
    m = re.search(r"Lead\s*gen.*?(\d[\d\s.,]*[KkMm]?)\s*book", txt, flags=re.I)
    if m: data["lead_gen_bookings"] = int(_normalize_number_fragment(m.group(1)))

    return data


# -------------------------------
# Bookings block
# -------------------------------

def parse_bookings_block(txt: str) -> Dict[str, Any]:
    """
    Parse 'BOOKINGS DETAILS' zone.
    """
    data = {
        "month_july_pct": None,
        "month_august_pct": None,
        "month_sept_pct": None,
        "top_dates_booked": None,
        "top_dates_searched": None,
        "top_parks_booked": None,
        "length_2n_pct": None,
        "length_3n_pct": None,
        "length_4n_pct": None,
        "raw_bookings_text": txt
    }

    # months
    # "July 46%, August 34% & September 7%"
    m = re.search(r"July\s*(\d+[.,]?\d*)%\D+August\s*(\d+[.,]?\d*)%\D+Sept(?:ember)?\s*(\d+[.,]?\d*)%", txt, flags=re.I)
    if m:
        data["month_july_pct"] = parse_percent(m.group(1)+"%")
        data["month_august_pct"] = parse_percent(m.group(2)+"%")
        data["month_sept_pct"] = parse_percent(m.group(3)+"%")

    # top dates booked
    m = re.search(r"Top\s*dates\s*booked\s*:\s*([^\n\r]+)", txt, flags=re.I)
    if m:
        data["top_dates_booked"] = _clean_csv_line(m.group(1))

    # top dates searched
    m = re.search(r"Top\s*dates\s*searched\s*:\s*([^\n\r]+)", txt, flags=re.I)
    if m:
        data["top_dates_searched"] = _clean_csv_line(m.group(1))

    # top parks
    # "BF 22%, BD 15% & LA 13%"
    m = re.search(r"Top\s*parks\s*booked\s*:\s*([^\n\r]+)", txt, flags=re.I)
    if m:
        parks_line = m.group(1)
        parts = re.findall(r"([A-Za-z]+)\s*(\d+[.,]?\d*)%", parks_line)
        if parts:
            data["top_parks_booked"] = ",".join(f"{code}:{parse_percent(num+'%'):.4f}" for code, num in parts)

    # lengths of stay
    m = re.search(r"Lengths?\s*of\s*stay\s*:\s*(.+)", txt, flags=re.I)
    if m:
        los_line = m.group(1)
        # 2 nights (33%), 3 nights (33%) & 4 nights (19%)
        parts = re.findall(r"(\d+)\s*night[s]?\s*\((\d+[.,]?\d*)%\)", los_line, flags=re.I)
        for n, pct in parts:
            val = parse_percent(pct+"%")
            if n == "2": data["length_2n_pct"] = val
            elif n == "3": data["length_3n_pct"] = val
            elif n == "4": data["length_4n_pct"] = val

    return data


def _clean_csv_line(line: str) -> str:
    # remove bullet separators like '&'
    line = line.strip().rstrip('.;')
    line = line.replace("&", ",")
    line = re.sub(r"\s+", " ", line)
    line = line.replace(" ,", ",")
    line = line.replace(" , ", ",")
    line = line.replace(" ,", ",")
    # unify comma separation
    line = line.replace(" ", "")
    # we want date tokens like Jul12? Actually we store as raw tokens; user can map.
    # Replace French months? - not needed for given english month abbreviations.
    # Reintroduce comma separation robustly:
    tokens = re.split(r"[;,]+", line)
    tokens = [t for t in re.split(r"[, ]+", line) if t]
    return ",".join(tokens)


# -------------------------------
# Main parse orchestrator
# -------------------------------

def parse_cpfr_slide(
    pptx_path: str,
    slide_number: Optional[int] = None,
    slide_title_contains: Optional[str] = "Sum up and main insights",
    week_start_date: Optional[str] = None
) -> Dict[str, Any]:
    """
    Extract structured data from the CPFR summary slide.

    slide_number: 1-based human index; if None we auto-detect by title substring.
    """
    prs = Presentation(pptx_path)
    slide = None

    if slide_number is not None:
        idx = slide_number - 1  # python-pptx zero-based
        if idx < 0 or idx >= len(prs.slides):
            raise ValueError(f"Slide number {slide_number} out of range (1..{len(prs.slides)})")
        slide = prs.slides[idx]
    else:
        # scan titles
        low = unidecode(slide_title_contains).lower()
        for s in prs.slides:
            title = ""
            if s.shapes.title:
                title = unidecode((s.shapes.title.text or "")).lower()
            if low in title:
                slide = s
                break
        if slide is None:
            raise RuntimeError("Slide not found by title.")

    # Build aggregated text for header region by concatenating all shapes above mid‑height?
    # Quick heuristic: shapes near top (top < 2in)
    header_texts = []
    overview_txt = ""
    offers_txt = ""
    bookings_txt = ""
    for sh in slide.shapes:
        if not hasattr(sh, "text"):
            continue
        t = sh.text.strip()
        if not t:
            continue
        norm = unidecode(t).lower()
        top = getattr(getattr(sh, "top", None), "emu", None)  # not needed
        # classify
        if "overview" in norm and "performances" in norm:
            overview_txt = t
        elif "focus" in norm and "offer" in norm:
            offers_txt = t
        elif "booking" in norm and "detail" in norm:
            bookings_txt = t
        else:
            # Likely header KPI shapes
            header_texts.append(t)

    kpi_header_txt = "\n".join(header_texts)

    kpi_data = parse_kpi_header(kpi_header_txt)
    ov_data = parse_overview_block(overview_txt)
    of_data = parse_offers_block(offers_txt)
    bk_data = parse_bookings_block(bookings_txt)

    # merge
    data = {
        "week_start_date": week_start_date,
        "weekly_summary": {**{k: v for k,v in kpi_data.items() if not k.startswith("vs_")}, **ov_data},
        "offers_focus": {k: v for k,v in of_data.items() if not k.startswith("raw_")},
        "bookings_details": {k: v for k,v in bk_data.items() if not k.startswith("raw_")},
        "variations": {k: v for k,v in kpi_data.items() if k.startswith("vs_")},
        "raw_text": {
            "header": kpi_header_txt,
            "overview": ov_data.get("raw_overview_text", overview_txt),
            "offers": of_data.get("raw_offers_text", offers_txt),
            "bookings": bk_data.get("raw_bookings_text", bookings_txt)
        }
    }

    # flatten variations into weekly_summary keys defined in schema
    data["weekly_summary"].update(data["variations"])

    return data


# -------------------------------
# CLI
# -------------------------------

def cli():
    ap = argparse.ArgumentParser(description="Parse CPFR slide into JSON.")
    ap.add_argument("pptx", help="Path to PPTX file.")
    ap.add_argument("--slide", type=int, default=31, help="Human slide number (1-based). Default=31.")
    ap.add_argument("--week-start", type=str, required=True, help="ISO date for week start (YYYY-MM-DD).")
    ap.add_argument("--out", type=str, default="-", help="Output JSON file or '-' for stdout.")
    args = ap.parse_args()

    data = parse_cpfr_slide(args.pptx, slide_number=args.slide, week_start_date=args.week_start)
    out_json = json.dumps(data, indent=2, ensure_ascii=False)

    if args.out == "-":
        print(out_json)
    else:
        Path(args.out).write_text(out_json, encoding="utf-8")

if __name__ == "__main__":
    cli() 