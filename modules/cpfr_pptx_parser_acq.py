"""
cpfr_pptx_parser_acq.py

Parsing 'CPFR LW – Acquisition Channel Analysis' slide (4 colonnes : SEA, SEO, OM, CRM).
Segmentation spatiale -> extraction texte -> parsing sémantique -> payload structuré.

Dépendances : python-pptx, unidecode, re.
"""

import re
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Tuple, List, Optional

from pptx import Presentation
from unidecode import unidecode

# ============================================================
# --------- Helpers: numeric parsing (shared logic) ----------
# ============================================================

def _clean_spaces(s: str) -> str:
    return re.sub(r"\s+", " ", s.strip())

def _normalize_number_fragment(raw: str) -> float:
    """
    Convert strings like '1,4M', '118K', '367K', '2 475', '0,53' to float.
    Currency / suffix € ignored; percent not divided here.
    """
    txt = raw.strip()
    txt = txt.replace("\xa0", "").replace(" ", "")
    # drop trailing euro
    if txt.endswith(("€", "E", "e")):
        txt = txt[:-1]
    mult = 1
    if txt.lower().endswith("k"):
        mult = 1_000
        txt = txt[:-1]
    elif txt.lower().endswith("m"):
        mult = 1_000_000
        txt = txt[:-1]
    txt = txt.replace(",", ".")
    try:
        return float(txt) * mult
    except ValueError:
        return 0.0

def parse_currency(raw: str) -> Optional[float]:
    m = re.search(r"([+-]?\d[\d\s.,]*[KkMm]?)(?:€|e)?", raw)
    if not m:
        return None
    return _normalize_number_fragment(m.group(1))

def parse_percent(raw: str) -> Optional[float]:
    """
    Convert '+50%' → 0.5; '-14%' → -0.14; '11,5%' ok; parentheses ignored.
    """
    txt = raw.strip()
    sign = 1
    if txt.startswith("+"):
        sign = 1; txt = txt[1:]
    elif txt.startswith("-"):
        sign = -1; txt = txt[1:]
    txt = txt.replace("%","").replace(" ","").replace(",",".")
    try:
        return sign * (float(txt)/100.0)
    except ValueError:
        return None

def parse_int(raw: str) -> Optional[int]:
    val = parse_currency(raw)
    if val is None:
        return None
    return int(round(val))

def parse_date_dmy_or_dmy_no_year(raw: str, ref_year: Optional[int] = None) -> Optional[str]:
    """
    Accept '15/07', '15-07', '15.07' -> iso date using ref_year (default: current year).
    """
    m = re.search(r"(\d{1,2})[./-](\d{1,2})(?:[./-](\d{2,4}))?", raw)
    if not m:
        return None
    d = int(m.group(1)); mth = int(m.group(2))
    y = m.group(3)
    if y:
        y = int(y)
        if y < 100: y += 2000
    else:
        y = ref_year or datetime.today().year
    try:
        return datetime(y, mth, d).strftime("%Y-%m-%d")
    except ValueError:
        return None


# ============================================================
# ------------ Slide segmentation (columns/rows) -------------
# ============================================================

def _assign_column(left_emu: int, width_emu: int, slide_width_emu: int) -> int:
    """
    Return column index 0..3 from shape center X.
    """
    center_x = left_emu + width_emu/2
    col_width = slide_width_emu / 4.0
    col = int(center_x // col_width)
    return min(max(col,0),3)

def _assign_band(top_emu: int, slide_height_emu: int, header_pct=0.2, footer_pct=0.8) -> int:
    """
    Return row band 0=header,1=body,2=footer.
    """
    top_ratio = top_emu / slide_height_emu
    if top_ratio < header_pct:
        return 0
    elif top_ratio > footer_pct:
        return 2
    return 1

def _collect_slide_text_by_grid(slide, header_pct=0.2, footer_pct=0.8) -> Dict[Tuple[int,int], List[str]]:
    """
    Collect raw text lines grouped by (col,band).
    """
    # simpler / stable:
    try:
        slide_width = slide.part.slide_width
        slide_height = slide.part.slide_height
    except AttributeError:
        # fallback to presentation object
        prs = slide.part.presentation_part.presentation
        slide_width = prs.slide_width
        slide_height = prs.slide_height

    buckets = {(c,b):[] for c in range(4) for b in range(3)}

    for sh in slide.shapes:
        if not hasattr(sh, "text") or not sh.text:
            continue
        txt = sh.text.strip()
        if not txt:
            continue
        col = _assign_column(sh.left, sh.width, slide_width)
        band = _assign_band(sh.top, slide_height, header_pct, footer_pct)
        buckets[(col,band)].append((sh.top, sh.left, txt))

    # sort & join
    out = {}
    for k, arr in buckets.items():
        if not arr:
            continue
        arr_sorted = sorted(arr, key=lambda r: (r[0], r[1]))
        txt_join = "\n".join(x[2] for x in arr_sorted)
        out[k] = txt_join
    return out


# ============================================================
# ---------------- Column content parsers --------------------
# ============================================================

# ---------- SEA ----------
def _parse_sea_block(body_txt: str) -> Dict[str, Any]:
    """
    Parse SEA stats body (WoW GA4 Sessions/Bookings/Revenue/Costs; CVR; Summer sales bullet bookings).
    """
    t = body_txt

    # WoW GA4 block: -7% Sessions, -18% Bookings, -14% Revenue, -10% Costs
    wow_sessions = _search_percent_after_label(t, r"Sessions")
    wow_bookings = _search_percent_after_label(t, r"Bookings")
    wow_revenue = _search_percent_after_label(t, r"Revenue")
    wow_costs = _search_percent_after_label(t, r"Costs")

    # CVR vs Last Week & vs LY
    cvr_vs_lw = _search_percent_after_label(t, r"CVR\s*vs\s*Last\s*Week")
    # Within line you also see "vs LY : +60%" ; we grab that
    cvr_vs_ly = _search_percent_after_label(t, r"vs\s*LY")

    # Summer Sales bullet bookings
    promo_ext = _search_int_after_label(t, r"Promo\s*Extension\s*:\s*(\d[\d\s.,KkMm]*)\s*Book")
    pmax_asset = _search_int_after_label(t, r"Pmax\s*Asset\s*:\s*(\d[\d\s.,KkMm]*)\s*Book")
    sitelink   = _search_int_after_label(t, r"Sitelink\s*:\s*(\d[\d\s.,KkMm]*)\s*Book")

    return {
        "wow_sessions": wow_sessions,
        "wow_bookings": wow_bookings,
        "wow_revenue": wow_revenue,
        "wow_costs": wow_costs,
        "cvr_vs_lw": cvr_vs_lw,
        "cvr_vs_ly": cvr_vs_ly,
        "campaign_notes": [
            {"campaign_name":"Promo Extension","metric_bookings":promo_ext},
            {"campaign_name":"Pmax Asset","metric_bookings":pmax_asset},
            {"campaign_name":"Sitelink","metric_bookings":sitelink}
        ]
    }

# ---------- SEO ----------
def _parse_seo_block(body_txt: str) -> Dict[str, Any]:
    t = body_txt

    # Traffic on Brand
    brand_impr  = _search_percent_after_label(t, r"Impressions:\s*([+-]?\d[\d.,]*)%\s*\(YoY\)", raw=True)
    brand_click = _search_percent_after_label(t, r"Clicks:\s*([+-]?\d[\d.,]*)%\s*\(YoY\)", raw=True, nth=1)
    brand_ctr   = _search_percent_after_label(t, r"CTR:\s*([+-]?\d[\d.,]*)%\s*\(YoY\)", raw=True, nth=2)
    brand_pos   = _search_float_after_label_generic(t, r"Average\s*Position:\s*([\d.,]+)")

    # Traffic on Non-Brand
    nb_impr  = _search_percent_after_label(t, r"Impressions:\s*([+-]?\d[\d.,]*)%\s*\(YoY\)", raw=True, start_after="Traffic on Non-Brand")
    nb_click = _search_percent_after_label(t, r"Clicks:\s*([+-]?\d[\d.,]*)%\s*\(YoY\)", raw=True, start_after="Traffic on Non-Brand", nth=1)
    nb_ctr   = _search_percent_after_label(t, r"CTR:\s*([+-]?\d[\d.,]*)%\s*\(YoY\)", raw=True, start_after="Traffic on Non-Brand", nth=2)
    nb_pos   = _search_float_after_label_generic(t, r"Average\s*Position:\s*([\d.,]+)", start_after="Traffic on Non-Brand")

    # Top branded / non branded / specific
    top_branded       = _list_after_label(t, r"Top\s*branded\s*request\s*:\s*(.+)")
    top_non_branded   = _list_after_label(t, r"Top\s*non\s*branded\s*request\s*:\s*(.+)")
    top_specific_brand= _list_after_label(t, r"Top\s*specific\s*brand\s*:\s*(.+)")

    return {
        "brand": {
            "impressions_yoy": brand_impr,
            "clicks_yoy": brand_click,
            "ctr_yoy": brand_ctr,
            "avg_position": brand_pos
        },
        "non_brand": {
            "impressions_yoy": nb_impr,
            "clicks_yoy": nb_click,
            "ctr_yoy": nb_ctr,
            "avg_position": nb_pos
        },
        "top_branded": top_branded,
        "top_non_branded": top_non_branded,
        "top_specific_brand": top_specific_brand,
        "campaign_notes": []
    }

# ---------- OM ----------
def _parse_om_block(body_txt: str) -> Dict[str, Any]:
    t = body_txt
    # headline: Traffic : +50% (WoW) // +74% (YoY)
    traffic_wow, traffic_yoy = _parse_dual_pct_line(t, r"Traffic\s*:\s*")
    trans_wow, trans_yoy     = _parse_dual_pct_line(t, r"Transaction\s*:\s*")
    rev_wow, rev_yoy         = _parse_dual_pct_line(t, r"Revenue\s*:\s*")

    # Affiliation line: Revenue -20% WoW / -44% YoY
    aff_rev_wow, aff_rev_yoy = _parse_dual_pct_line(t, r"Affiliation\s*:\s*Revenue\s*")
    # R-Advertising revenue +68% WoW
    radv_rev_wow = _search_percent_after_label(t, r"R-Advertising.*?\+?(-?\d[\d.,]*)%", allow_neg=True)
    # Retargeting revenue +73% WoW / +145% YoY
    ret_rev_wow, ret_rev_yoy = _parse_dual_pct_line(t, r"Retargeting\s*:\s*Revenue\s*")
    # SMP Sessions +25% WoW / +154% YoY
    smp_ses_wow, smp_ses_yoy = _parse_dual_pct_line(t, r"SMP.*?Sessions\s*")
    # Display + Native Sessions +16% WoW // -8% YoY
    dn_ses_wow, dn_ses_yoy   = _parse_dual_pct_line(t, r"Display\s*\+\s*Native.*?Sessions\s*")

    # Build notes
    camp_notes = []
    if aff_rev_wow is not None or aff_rev_yoy is not None:
        camp_notes.append({"campaign_name":"Affiliation","note":"End of coupon & cashback limit"})
    if radv_rev_wow is not None:
        camp_notes.append({"campaign_name":"R-Advertising","note":"Flash sale & LM activations"})
    if ret_rev_wow is not None or ret_rev_yoy is not None:
        camp_notes.append({"campaign_name":"Retargeting","note":"Push flash sale"})
    if smp_ses_wow is not None:
        camp_notes.append({"campaign_name":"SMP","note":"EB + LM offers"})
    if dn_ses_wow is not None:
        camp_notes.append({"campaign_name":"Display+Native","note":"Good performance"})

    return {
        "wow_sessions": traffic_wow,
        "yoy_sessions": traffic_yoy,
        "wow_bookings": trans_wow,
        "yoy_bookings": trans_yoy,
        "wow_revenue": rev_wow,
        "yoy_revenue": rev_yoy,
        "affiliation_revenue_wow": aff_rev_wow,
        "affiliation_revenue_yoy": aff_rev_yoy,
        "retargeting_revenue_wow": ret_rev_wow,
        "retargeting_revenue_yoy": ret_rev_yoy,
        "smp_sessions_wow": smp_ses_wow,
        "smp_sessions_yoy": smp_ses_yoy,
        "displaynative_sessions_wow": dn_ses_wow,
        "displaynative_sessions_yoy": dn_ses_yoy,
        "campaign_notes": camp_notes
    }

# ---------- CRM ----------
def _parse_crm_block(body_txt: str) -> Dict[str, Any]:
    t = body_txt
    # General: vs LY : +23% visits, +7% bookings, +15% revenue
    gen_ly_vis, gen_ly_book, gen_ly_rev = _parse_triple_pct_line(t, r"General:\s*vs\s*LY\s*:\s*")
    # vs LW : +22% visits, +33% bookings, +40% revenue
    gen_lw_vis, gen_lw_book, gen_lw_rev = _parse_triple_pct_line(t, r"vs\s*LW\s*:\s*")
    # Tactical Last Week: Booking 115, Turnover 118k €
    last_book = _search_int_after_label(t, r"Booking\s*:\s*([\d\s.,KkMm]+)")
    last_turn = _search_currency_after_label(t, r"Turnover\s*:\s*([\d\s.,KkMm]+)")
    # Strategic JU25: Booking +32.4% vs LY; NBR +16.8% vs LY; Incremental : 526K€
    ju25_book = _search_percent_after_label(t, r"Booking\s*:\s*([+-]?\d[\d.,]*)%\s*vs\s*LY")
    ju25_nbr  = _search_percent_after_label(t, r"NBR\s*:\s*([+-]?\d[\d.,]*)%\s*vs\s*LY")
    ju25_incr = _search_currency_after_label(t, r"Incremental\s*:\s*([\d\s.,KkMm]+)")

    # This week actions (just keep raw)
    # B2C / B2B bullet detection -> notes
    b2c_flag = bool(re.search(r"B2C\s*:\s*Reminder\s*Summer\s*flash\s*sales", t, flags=re.I))
    b2b_flag = bool(re.search(r"B2B\s*:\s*", t, flags=re.I))

    camp_notes = []
    if last_book is not None or last_turn is not None:
        camp_notes.append({"campaign_name":"CRM Summer flash sales 2","metric_bookings":last_book,"metric_revenue":last_turn})
    if b2c_flag:
        camp_notes.append({"campaign_name":"CRM B2C Reminder Flash","note":"Reminder Summer flash sales"})
    if b2b_flag:
        camp_notes.append({"campaign_name":"CRM B2B Septembre","note":"Petits prix septembre"})

    return {
        "yoy_sessions": gen_ly_vis,
        "yoy_bookings": gen_ly_book,
        "yoy_revenue": gen_ly_rev,
        "wow_sessions": gen_lw_vis,
        "wow_bookings": gen_lw_book,
        "wow_revenue": gen_lw_rev,
        "tactical_last_bookings": last_book,
        "tactical_last_revenue": last_turn,
        "strategic_booking_yoy": ju25_book,
        "strategic_nbr_yoy": ju25_nbr,
        "strategic_incremental": ju25_incr,
        "campaign_notes": camp_notes
    }


# ============================================================
# --------------- Regex utility sub-parsers ------------------
# ============================================================

def _search_percent_after_label(text: str, label_pattern: str, raw=False, start_after=None, nth=0, allow_neg=True) -> Optional[float]:
    """
    Generic "label ... +X%" extraction.
    If raw=True label_pattern is full regex w/ capture group -> we parse group directly.
    start_after restricts search after first match of that token.
    nth selects nth occurrence if multiple.
    """
    search_space = text
    if start_after:
        idx = re.search(start_after, text, flags=re.I)
        if idx:
            search_space = text[idx.end():]
    if raw:
        matches = re.findall(label_pattern, search_space, flags=re.I)
        if not matches or nth >= len(matches):
            return None
        return parse_percent(matches[nth] + "%")
    # simple case: find label then % following
    pat = re.compile(label_pattern + r".*?([+-]?\d[\d.,]*)%", flags=re.I|re.S)
    matches = pat.findall(search_space)
    if not matches or nth >= len(matches):
        return None
    return parse_percent(matches[nth] + "%")

def _search_int_after_label(text: str, regex_pattern: str) -> Optional[int]:
    m = re.search(regex_pattern, text, flags=re.I)
    if not m:
        return None
    return parse_int(m.group(1))

def _search_currency_after_label(text: str, regex_pattern: str) -> Optional[float]:
    m = re.search(regex_pattern, text, flags=re.I)
    if not m:
        return None
    return parse_currency(m.group(1))

def _search_float_after_label_generic(text: str, regex_pattern: str, start_after=None) -> Optional[float]:
    search_space = text
    if start_after:
        idx = re.search(start_after, text, flags=re.I)
        if idx:
            search_space = text[idx.end():]
    m = re.search(regex_pattern, search_space, flags=re.I)
    if not m:
        return None
    try:
        return float(m.group(1).replace(",", "."))
    except ValueError:
        return None

def _list_after_label(text: str, regex_pattern: str) -> List[str]:
    m = re.search(regex_pattern, text, flags=re.I)
    if not m:
        return []
    line = m.group(1).strip()
    line = line.strip('"')
    # split on comma or semicolon
    parts = re.split(r"[;,]", line)
    parts = [p.strip(' "').strip("'") for p in parts if p.strip()]
    return parts

def _parse_dual_pct_line(text: str, prefix_pattern: str) -> Tuple[Optional[float], Optional[float]]:
    """
    Parse lines like 'Traffic : +50% (WoW) // +74% (YoY)'.
    Returns (wow, yoy).
    """
    m = re.search(prefix_pattern + r"([+-]?\d[\d.,]*)%\s*\(WoW\).*?([+-]?\d[\d.,]*)%\s*\(YoY\)", text, flags=re.I|re.S)
    if not m:
        return None, None
    return parse_percent(m.group(1)+"%"), parse_percent(m.group(2)+"%")

def _parse_triple_pct_line(text: str, prefix_pattern: str) -> Tuple[Optional[float], Optional[float], Optional[float]]:
    """
    Parse 'vs LY : +23% visits, +7% bookings, +15% revenue'.
    Returns tuple (visits, bookings, revenue).
    """
    m = re.search(prefix_pattern + r"([+-]?\d[\d.,]*)%\s*visits[^%]*?([+-]?\d[\d.,]*)%\s*bookings[^%]*?([+-]?\d[\d.,]*)%\s*revenue", text, flags=re.I)
    if not m:
        return None, None, None
    return parse_percent(m.group(1)+"%"), parse_percent(m.group(2)+"%"), parse_percent(m.group(3)+"%")


# ============================================================
# ---------- Main parse function for Acquisition slide -------
# ============================================================

def parse_acquisition_slide(
    pptx_path: str,
    slide_number: int = 32,
    week_start_date: Optional[str] = None,
    header_pct: float = 0.2,
    footer_pct: float = 0.8
) -> Dict[str, Any]:
    """
    Parse the Acquisition Channel Analysis slide.
    Returns structured dict w/ SEA, SEO, OM, CRM blocks + last_update dates.
    """
    prs = Presentation(pptx_path)
    idx = slide_number - 1
    if idx < 0 or idx >= len(prs.slides):
        raise ValueError(f"Slide {slide_number} out of bounds.")
    slide = prs.slides[idx]

    buckets = _collect_slide_text_by_grid(slide, header_pct=header_pct, footer_pct=footer_pct)

    # Build column text (header/body/footer)
    cols = {}
    col_names = ["SEA","SEO","OM","CRM"]
    for ci, cname in enumerate(col_names):
        header_txt = buckets.get((ci,0), "")
        body_txt   = buckets.get((ci,1), "")
        footer_txt = buckets.get((ci,2), "")
        cols[cname] = {"header": header_txt, "body": body_txt, "footer": footer_txt}

    # Parse each column
    sea_metrics = _parse_sea_block(cols["SEA"]["body"])
    seo_metrics = _parse_seo_block(cols["SEO"]["body"])
    om_metrics  = _parse_om_block(cols["OM"]["body"])
    crm_metrics = _parse_crm_block(cols["CRM"]["body"])

    # parse last updates (footer lines)
    for cname in col_names:
        cols[cname]["last_update"] = _parse_last_update(cols[cname]["footer"])

    sea_metrics["last_update"] = cols["SEA"]["last_update"]
    seo_metrics["last_update"] = cols["SEO"]["last_update"]
    om_metrics["last_update"]  = cols["OM"]["last_update"]
    crm_metrics["last_update"] = cols["CRM"]["last_update"]

    # raw text attachments
    sea_metrics["raw"] = cols["SEA"]["body"]
    seo_metrics["raw"] = cols["SEO"]["body"]
    om_metrics["raw"]  = cols["OM"]["body"]
    crm_metrics["raw"] = cols["CRM"]["body"]

    return {
        "week_start_date": week_start_date,
        "acquisition": {
            "SEA": sea_metrics,
            "SEO": seo_metrics,
            "OM": om_metrics,
            "CRM": crm_metrics
        },
        "raw_columns": cols
    }


def _parse_last_update(txt: str) -> Optional[str]:
    """
    Parse 'Last update SEA : 15/07' style -> ISO date (assume current year).
    """
    if not txt:
        return None
    return parse_date_dmy_or_dmy_no_year(txt)


# ============================================================
# -------- Mapping parsed metrics -> DB payload tables --------
# ============================================================

def build_acquisition_db_payload(parsed: Dict[str, Any]) -> Dict[str, Any]:
    """
    Convert parse_acquisition_slide() output into DB insert payload
    matching acquisition_channels, channel_campaign_notes, channel_seo_detail.
    """
    week_start = parsed.get("week_start_date")
    acq = parsed["acquisition"]

    # Build acquisition_channels rows
    acq_rows = []
    camp_notes_rows = []
    seo_detail_rows = []

    # SEA
    sea = acq["SEA"]
    acq_rows.append({
        "channel_code": "SEA",
        "wow_sessions": sea.get("wow_sessions"),
        "wow_bookings": sea.get("wow_bookings"),
        "wow_revenue": sea.get("wow_revenue"),
        "wow_costs": sea.get("wow_costs"),
        "cvr_vs_lw": sea.get("cvr_vs_lw"),
        "cvr_vs_ly": sea.get("cvr_vs_ly"),
        "comments": sea.get("raw")
    })
    camp_notes_rows += _mk_notes("SEA", sea.get("campaign_notes", []))

    # SEO
    seo = acq["SEO"]
    acq_rows.append({
        "channel_code": "SEO",
        "wow_sessions": None,
        "wow_bookings": None,
        "wow_revenue": None,
        "wow_costs": None,
        "cvr_vs_lw": None,
        "cvr_vs_ly": None,
        "comments": seo.get("raw")
    })
    camp_notes_rows += _mk_notes("SEO", seo.get("campaign_notes", []))
    # SEO detail table
    brand = seo.get("brand",{})
    nonb  = seo.get("non_brand",{})
    seo_detail_rows.append({
        "segment": "brand",
        "impressions_yoy": brand.get("impressions_yoy"),
        "clicks_yoy": brand.get("clicks_yoy"),
        "ctr_yoy": brand.get("ctr_yoy"),
        "avg_position": brand.get("avg_position")
    })
    seo_detail_rows.append({
        "segment": "non_brand",
        "impressions_yoy": nonb.get("impressions_yoy"),
        "clicks_yoy": nonb.get("clicks_yoy"),
        "ctr_yoy": nonb.get("ctr_yoy"),
        "avg_position": nonb.get("avg_position")
    })

    # OM
    om = acq["OM"]
    acq_rows.append({
        "channel_code": "OM",
        "wow_sessions": om.get("wow_sessions"),
        "yoy_sessions": om.get("yoy_sessions"),
        "wow_bookings": om.get("wow_bookings"),
        "yoy_bookings": om.get("yoy_bookings"),
        "wow_revenue": om.get("wow_revenue"),
        "yoy_revenue": om.get("yoy_revenue"),
        "comments": om.get("raw")
    })
    camp_notes_rows += _mk_notes("OM", om.get("campaign_notes", []))

    # CRM
    crm = acq["CRM"]
    acq_rows.append({
        "channel_code": "CRM",
        "wow_sessions": crm.get("wow_sessions"),
        "yoy_sessions": crm.get("yoy_sessions"),
        "wow_bookings": crm.get("wow_bookings"),
        "yoy_bookings": crm.get("yoy_bookings"),
        "wow_revenue": crm.get("wow_revenue"),
        "yoy_revenue": crm.get("yoy_revenue"),
        "comments": crm.get("raw")
    })
    camp_notes_rows += _mk_notes("CRM", crm.get("campaign_notes", []))

    return {
        "week_start_date": week_start,
        "acquisition_channels": acq_rows,
        "channel_campaign_notes": camp_notes_rows,
        "channel_seo_detail": seo_detail_rows
    }


def _mk_notes(channel_code: str, notes_list: List[Dict[str,Any]]) -> List[Dict[str,Any]]:
    out = []
    for n in notes_list:
        row = {
            "channel_code": channel_code,
            "campaign_name": n.get("campaign_name"),
            "metric_bookings": n.get("metric_bookings"),
            "metric_revenue": n.get("metric_revenue"),
            "note": n.get("note")
        }
        out.append(row)
    return out


# ============================================================
# -------------------------- CLI -----------------------------
# ============================================================

if __name__ == "__main__":
    import argparse, json
    ap = argparse.ArgumentParser(description="Parse CPFR Acquisition slide (SEA/SEO/OM/CRM).")
    ap.add_argument("pptx", help="Path to PPTX.")
    ap.add_argument("--slide", type=int, default=32, help="Human slide number (default=32).")
    ap.add_argument("--week-start", required=True, help="Week start date YYYY-MM-DD.")
    ap.add_argument("--out-json", default="-", help="Write parsed JSON to file or '-' for stdout.")
    ap.add_argument("--db-payload", action="store_true", help="Output DB payload mapping instead of raw parse.")
    args = ap.parse_args()

    parsed = parse_acquisition_slide(args.pptx, slide_number=args.slide, week_start_date=args.week_start)
    if args.db_payload:
        payload = build_acquisition_db_payload(parsed)
        out_obj = payload
    else:
        out_obj = parsed

    out_str = json.dumps(out_obj, indent=2, ensure_ascii=False)
    if args.out_json == "-":
        print(out_str)
    else:
        Path(args.out_json).write_text(out_str, encoding="utf-8") 