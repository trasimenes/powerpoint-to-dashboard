from pptx import Presentation
import re


def clean_text(text):
    """Nettoie et normalise le texte extrait"""
    if not text:
        return ""
    # Supprime les espaces multiples et les retours à la ligne
    text = re.sub(r'\s+', ' ', text.strip())
    return text


def parse_slide_text(slide):
    """Extrait tous les textes d'une slide avec une meilleure organisation"""
    texts = []
    
    print(f"[DEBUG] Parsing slide text - {len(slide.shapes)} formes trouvées")
    
    for i, shape in enumerate(slide.shapes):
        try:
            shape_type = type(shape).__name__
            print(f"[DEBUG] Forme {i}: {shape_type}")
            
            if hasattr(shape, "text") and shape.text:
                # Vérification que le texte est bien une chaîne
                if isinstance(shape.text, str):
                    text = clean_text(shape.text)
                    if text and len(text) > 2:  # Ignore les textes trop courts
                        texts.append(text)
                        print(f"[DEBUG] Texte extrait forme {i}: '{text}'")
                    else:
                        print(f"[DEBUG] Texte trop court ignoré forme {i}: '{text}'")
                else:
                    print(f"[DEBUG] Texte non-string forme {i}: {type(shape.text)}")
            else:
                print(f"[DEBUG] Forme {i} n'a pas de texte")
        except Exception as e:
            # Ignore les erreurs sur les formes problématiques
            print(f"[DEBUG] Erreur lors de l'extraction du texte forme {i}: {e}")
            continue
    
    # Trie par ordre d'apparition et supprime les doublons
    unique_texts = []
    for text in texts:
        if text not in unique_texts:
            unique_texts.append(text)
    
    print(f"[DEBUG] Textes finaux extraits: {unique_texts}")
    return unique_texts


def parse_table(slide):
    """Extrait les données de tableau avec une meilleure structure"""
    debug_info = {
        "total_shapes": len(slide.shapes),
        "table_shapes_found": 0,
        "tables_with_rows": 0,
        "tables_with_valid_data": 0,
        "shape_types": [],
        "errors": []
    }
    
    print(f"[DEBUG] Analyse de la slide - Nombre total de formes: {len(slide.shapes)}")
    
    for i, shape in enumerate(slide.shapes):
        try:
            shape_type = type(shape).__name__
            debug_info["shape_types"].append(shape_type)
            print(f"[DEBUG] Forme {i}: {shape_type}")
            
            if hasattr(shape, 'has_table'):
                print(f"[DEBUG] Forme {i} a l'attribut has_table: {shape.has_table}")
                
                if shape.has_table:
                    debug_info["table_shapes_found"] += 1
                    print(f"[DEBUG] Tableau trouvé dans la forme {i}")
                    
                    try:
                        table = shape.table
                        print(f"[DEBUG] Objet table créé avec succès")
                    except Exception as e:
                        print(f"[DEBUG] Erreur création objet table: {e}")
                        debug_info["errors"].append(f"Erreur création table forme {i}: {str(e)}")
                        continue
                    
                    try:
                        rows_count = len(table.rows)
                        print(f"[DEBUG] Tableau - Nombre de lignes: {rows_count}")
                    except Exception as e:
                        print(f"[DEBUG] Erreur accès table.rows: {e}")
                        debug_info["errors"].append(f"Erreur accès rows forme {i}: {str(e)}")
                        continue
                    
                    try:
                        cols_count = len(table.columns)
                        print(f"[DEBUG] Tableau - Nombre de colonnes: {cols_count}")
                    except Exception as e:
                        print(f"[DEBUG] Erreur accès table.columns: {e}")
                        debug_info["errors"].append(f"Erreur accès columns forme {i}: {str(e)}")
                        continue
                    
                    if not table.rows:
                        print(f"[DEBUG] Tableau sans lignes - ignoré")
                        continue
                    
                    debug_info["tables_with_rows"] += 1
                        
                    # Extraction des en-têtes
                    headers = []
                    try:
                        first_row = table.rows[0]
                        print(f"[DEBUG] Première ligne obtenue, {len(first_row.cells)} cellules")
                        
                        for j, cell in enumerate(first_row.cells):
                            try:
                                print(f"[DEBUG] Traitement en-tête cellule {j}")
                                cell_text = cell.text if hasattr(cell, 'text') else ""
                                print(f"[DEBUG] En-tête {j} - Texte brut: '{cell_text}' (type: {type(cell_text)})")
                                
                                if isinstance(cell_text, str):
                                    header_text = clean_text(cell_text)
                                else:
                                    header_text = ""
                                    
                                headers.append(header_text if header_text else f"Colonne {len(headers) + 1}")
                                print(f"[DEBUG] En-tête {j}: '{header_text}'")
                            except Exception as e:
                                headers.append(f"Colonne {len(headers) + 1}")
                                print(f"[DEBUG] Erreur en-tête {j}: {e}")
                                debug_info["errors"].append(f"Erreur en-tête {j}: {str(e)}")
                    except Exception as e:
                        print(f"[DEBUG] Erreur accès première ligne: {e}")
                        debug_info["errors"].append(f"Erreur accès première ligne: {str(e)}")
                        continue
                    
                    # Extraction des données
                    rows = []
                    try:
                        total_data_rows = len(table.rows) - 1
                        print(f"[DEBUG] Traitement des lignes de données (total: {total_data_rows})")
                        
                        if total_data_rows <= 0:
                            print(f"[DEBUG] Pas de lignes de données (seulement en-têtes)")
                            continue
                    except Exception as e:
                        print(f"[DEBUG] Erreur calcul nombre lignes données: {e}")
                        debug_info["errors"].append(f"Erreur calcul lignes données: {str(e)}")
                        continue
                    
                    # Itération directe sur les lignes sans slice
                    for row_idx in range(1, len(table.rows)):
                        try:
                            row = table.rows[row_idx]
                            row_data = []
                            print(f"[DEBUG] Traitement ligne {row_idx} avec {len(row.cells)} cellules")
                            
                            for cell_idx, cell in enumerate(row.cells):
                                try:
                                    raw_text = cell.text if hasattr(cell, 'text') else ""
                                    print(f"[DEBUG] Ligne {row_idx}, Cellule {cell_idx} - Texte brut: '{raw_text}' (type: {type(raw_text)})")
                                    
                                    if isinstance(raw_text, str):
                                        cell_text = clean_text(raw_text)
                                        print(f"[DEBUG] Ligne {row_idx}, Cellule {cell_idx} - Texte nettoyé: '{cell_text}'")
                                    else:
                                        cell_text = ""
                                        print(f"[DEBUG] Ligne {row_idx}, Cellule {cell_idx} - Texte non-string converti en vide")
                                    
                                    row_data.append(cell_text if cell_text else "")
                                except Exception as e:
                                    row_data.append("")
                                    print(f"[DEBUG] Erreur cellule {row_idx},{cell_idx}: {e}")
                            
                            print(f"[DEBUG] Ligne {row_idx} - Données complètes: {row_data}")
                            
                            # Ne garde que les lignes qui ont au moins une donnée
                            has_valid_data = False
                            for cell in row_data:
                                if isinstance(cell, str) and cell.strip():
                                    has_valid_data = True
                                    print(f"[DEBUG] Ligne {row_idx} - Données valides trouvées: '{cell.strip()}'")
                                    break
                            
                            if has_valid_data:
                                rows.append(row_data)
                                print(f"[DEBUG] Ligne {row_idx} ajoutée au résultat")
                            else:
                                print(f"[DEBUG] Ligne {row_idx} ignorée (toutes cellules vides)")
                        except Exception as e:
                            # Ignore les lignes problématiques
                            print(f"[DEBUG] Erreur ligne {row_idx}: {e}")
                            debug_info["errors"].append(f"Erreur ligne {row_idx}: {str(e)}")
                            continue
                    
                    if rows:
                        debug_info["tables_with_valid_data"] += 1
                        print(f"[DEBUG] Tableau valide trouvé avec {len(rows)} lignes de données")
                        
                        return {
                            "headers": headers,
                            "rows": rows,
                            "total_rows": len(rows),
                            "total_columns": len(headers),
                            "debug": debug_info
                        }
                    else:
                        print(f"[DEBUG] Tableau trouvé mais sans données valides")
            else:
                print(f"[DEBUG] Forme {i} n'a pas d'attribut has_table")
        except Exception as e:
            error_msg = f"Erreur forme {i}: {str(e)}"
            debug_info["errors"].append(error_msg)
            print(f"[DEBUG] {error_msg}")
            continue
    
    print(f"[DEBUG] Résumé: {debug_info['table_shapes_found']} tableaux trouvés, {debug_info['tables_with_valid_data']} avec données valides")
    
    return {
        "headers": [],
        "rows": [],
        "total_rows": 0,
        "total_columns": 0,
        "debug": debug_info
    }


def extract_kpis_from_text(texts):
    """Extrait et organise les KPIs à partir des textes"""
    kpis = []
    
    for text in texts:
        # Vérification que text est bien une chaîne de caractères
        if not isinstance(text, str):
            continue
            
        # Détecte les patterns de KPIs (chiffres, pourcentages, etc.)
        if re.search(r'\d+[.,]?\d*%?', text):
            kpis.append(text)
        # Détecte les textes courts qui pourraient être des KPIs
        elif len(text) < 100 and any(keyword in text.lower() for keyword in 
                ['kpi', 'indicateur', 'performance', 'résultat', 'objectif', 'cible']):
            kpis.append(text)
        # Ajoute les autres textes importants
        elif len(text) > 10:
            kpis.append(text)
    
    return kpis


def extract_cpfr_data_from_slide31(slide):
    """
    Extrait spécifiquement les données CPFR de la slide 31
    Format attendu: KPI principal + variations LY/LW
    Ex: "342K sessions +6% vs LY, -4% vs LW"
    """
    cpfr_data = {
        'sessions': None,
        'revenue_b2c': None,
        'average_basket_value': None,
        'conversion_rate': None,
        'nb_bookings': None,
        'best_day': None,
        'best_day_sessions': None,
        'best_day_revenue': None,
        # Variations vs Last Year
        'vs_ly_sessions': None,
        'vs_ly_revenue': None,
        'vs_ly_abv': None,
        'vs_ly_cr': None,
        'vs_ly_bookings': None,
        # Variations vs Last Week
        'vs_lw_sessions': None,
        'vs_lw_revenue': None,
        'vs_lw_abv': None,
        'vs_lw_cr': None,
        'vs_lw_bookings': None,
        # Autres données
        'last_minute_pct': None,
        'early_booking_pct': None,
        'month_july_pct': None,
        'month_august_pct': None,
        'month_sept_pct': None,
        'top_dates_booked': [],
        'top_dates_searched': [],
        'top_parks': [],
        'lengths_of_stay': [],
        'insights': []
    }
    
    # Extraction de tous les textes de la slide
    texts = parse_slide_text(slide)
    
    print(f"[DEBUG] Début extraction CPFR slide 31")
    
    for i, text in enumerate(texts):
        text_lower = text.lower()
        print(f"[DEBUG] Analyse texte {i}: '{text}'")
        
        # SESSIONS avec variations - Pattern: "342K sessions +6% vs LY, -4% vs LW"
        if 'session' in text_lower:
            # Extraire la valeur principale
            session_match = re.search(r'(\d+(?:[\s,\.]\d+)*)\s*k.*session', text_lower)
            if session_match:
                try:
                    value = session_match.group(1).replace(' ', '').replace(',', '').replace('.', '')
                    if value.isdigit():
                        cpfr_data['sessions'] = int(value) * 1000
                        print(f"[DEBUG] Sessions trouvées: {cpfr_data['sessions']}")
                except:
                    pass
            
            # Extraire les variations LY/LW avec patterns plus précis
            ly_match = re.search(r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+(?:vs\s+)?ly', text_lower)
            if ly_match:
                try:
                    cpfr_data['vs_ly_sessions'] = float(ly_match.group(1).replace(',', '.')) / 100
                    print(f"[DEBUG] Sessions vs LY: {cpfr_data['vs_ly_sessions']}")
                except:
                    pass
            
            lw_match = re.search(r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+(?:vs\s+)?lw', text_lower)
            if lw_match:
                try:
                    cpfr_data['vs_lw_sessions'] = float(lw_match.group(1).replace(',', '.')) / 100
                    print(f"[DEBUG] Sessions vs LW: {cpfr_data['vs_lw_sessions']}")
                except:
                    pass
        
        # REVENUE - traiter séparément la valeur et les variations
        # Pattern pour la valeur: "2,27M€"
        if ('€' in text_lower and 'm' in text_lower) and not 'revenue' in text_lower:
            print(f"[DEBUG] Texte contenant €M (valeur): '{text}'")
            
            # Extraire la valeur principale - patterns simples pour capturer 2,27M
            revenue_patterns = [
                r'(\d+(?:[,\.]\d+)*)\s*m\s*€',              # "2,27M €"
                r'(\d+(?:[,\.]\d+)*)\s*m€',                 # "2,27M€"
                r'€\s*(\d+(?:[,\.]\d+)*)\s*m',              # "€ 2,27M"
            ]
            
            for pattern in revenue_patterns:
                revenue_match = re.search(pattern, text_lower)
                if revenue_match:
                    try:
                        value = revenue_match.group(1).replace(',', '.')
                        if value.replace('.', '').isdigit():
                            new_value = float(value) * 1000000
                            print(f"[DEBUG] Revenue CANDIDAT: {new_value} (pattern: {pattern}, value: {value})")
                            # Vérifier si c'est un écrasement
                            if 'revenue_b2c' in cpfr_data:
                                print(f"[DEBUG] ATTENTION: Ecrasement revenue_b2c {cpfr_data['revenue_b2c']} par {new_value}")
                            cpfr_data['revenue_b2c'] = new_value
                            print(f"[DEBUG] Revenue FINAL: {cpfr_data['revenue_b2c']}")
                            break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing revenue: {e}")
                        continue
            
        # REVENUE - traiter les variations LY/LW: "Web B2C Global revenue +11% VS LY -12% VS LW"
        if 'revenue' in text_lower:
            print(f"[DEBUG] Texte contenant revenue (variations): '{text}'")
            # Pattern LY plus flexible - capturer le % directement avant "VS LY"
            ly_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+ly',      # "+11% VS LY"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+ly',           # "+11% LY"
                r'vs\s+ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LY: +11%"
                r'ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LY +11%"
            ]
            
            for pattern in ly_patterns:
                ly_match = re.search(pattern, text_lower)
                if ly_match:
                    try:
                        cpfr_data['vs_ly_revenue'] = float(ly_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] Revenue vs LY: {cpfr_data['vs_ly_revenue']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LY: {e}")
                        continue
            
            # Pattern LW plus flexible - capturer le % directement avant "VS LW"
            lw_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+lw',      # "-12% VS LW"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+lw',           # "-12% LW"
                r'vs\s+lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LW: -12%"
                r'lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LW -12%"
            ]
            
            for pattern in lw_patterns:
                lw_match = re.search(pattern, text_lower)
                if lw_match:
                    try:
                        cpfr_data['vs_lw_revenue'] = float(lw_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] Revenue vs LW: {cpfr_data['vs_lw_revenue']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LW: {e}")
                        continue
        
        # AVERAGE BASKET VALUE - traiter séparément la valeur et les variations
        # Pattern pour la valeur: "917€"
        if '€' in text_lower and not 'basket' in text_lower and not 'panier' in text_lower and not 'm' in text_lower:
            print(f"[DEBUG] Texte contenant € (valeur ABV): '{text}'")
            
            # Extraire la valeur principale - patterns pour capturer 917€
            abv_patterns = [
                r'(\d+(?:[\s,\.]\d+)*)\s*€',                         # "917€"
                r'€\s*(\d+(?:[\s,\.]\d+)*)',                         # "€ 917"
            ]
            
            for pattern in abv_patterns:
                abv_match = re.search(pattern, text_lower)
                if abv_match:
                    try:
                        value = abv_match.group(1).replace(' ', '').replace(',', '.')
                        if value.replace('.', '').isdigit():
                            # Vérifier que c'est une valeur cohérente pour un panier moyen (100-2000€)
                            amount = float(value)
                            if 100 <= amount <= 2000:
                                cpfr_data['average_basket_value'] = amount
                                print(f"[DEBUG] ABV trouvé: {cpfr_data['average_basket_value']} (pattern: {pattern})")
                                break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing ABV: {e}")
                        continue
            
        # AVERAGE BASKET VALUE - traiter les variations: "Average basket value -15% VS LY +8% VS LW"
        if 'basket' in text_lower or 'panier' in text_lower:
            print(f"[DEBUG] Texte contenant basket/panier (variations): '{text}'")
            # Extraire les variations LY/LW avec patterns plus robustes
            # Pattern LY plus flexible - capturer le % directement avant "VS LY"
            ly_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+ly',      # "-15% VS LY"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+ly',           # "-15% LY"
                r'vs\s+ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LY: -15%"
                r'ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LY -15%"
            ]
            
            for pattern in ly_patterns:
                ly_match = re.search(pattern, text_lower)
                if ly_match:
                    try:
                        cpfr_data['vs_ly_abv'] = float(ly_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] ABV vs LY: {cpfr_data['vs_ly_abv']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LY: {e}")
                        continue
            
            # Pattern LW plus flexible - capturer le % directement avant "VS LW"  
            lw_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+lw',      # "+8% VS LW"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+lw',           # "+8% LW"
                r'vs\s+lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LW: +8%"
                r'lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LW +8%"
            ]
            
            for pattern in lw_patterns:
                lw_match = re.search(pattern, text_lower)
                if lw_match:
                    try:
                        cpfr_data['vs_lw_abv'] = float(lw_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] ABV vs LW: {cpfr_data['vs_lw_abv']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LW: {e}")
                        continue
        
        # CONVERSION RATE - traiter séparément la valeur et les variations
        # Pattern pour la valeur: "0,53%"
        if '%' in text and not 'conversion' in text_lower and not 'taux' in text_lower and not 'basket' in text_lower and not 'panier' in text_lower and not 'session' in text_lower and not 'booking' in text_lower and not 'revenue' in text_lower:
            print(f"[DEBUG] Texte contenant % (valeur conversion): '{text}'")
            
            # Extraire la valeur principale - patterns pour capturer 0,53%
            conv_patterns = [
                r'(\d+[,\.]\d+)\s*%',                               # "0,53%"
                r'(\d+[,\.]\d+)\s*%\s*$',                           # "0,53%" en fin de ligne
            ]
            
            for pattern in conv_patterns:
                conv_match = re.search(pattern, text_lower)
                if conv_match:
                    try:
                        value = conv_match.group(1).replace(',', '.')
                        if value.replace('.', '').isdigit():
                            # Vérifier que c'est une valeur cohérente pour un taux de conversion (0.1% à 5%)
                            amount = float(value)
                            if 0.1 <= amount <= 5.0:
                                cpfr_data['conversion_rate'] = amount / 100
                                print(f"[DEBUG] Conversion Rate trouvé: {cpfr_data['conversion_rate']} (pattern: {pattern})")
                                break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing conversion: {e}")
                        continue
            
        # CONVERSION RATE - traiter les variations: "Conversion rate +12% VS LY -14% VS LW"
        if 'conversion' in text_lower or 'taux' in text_lower:
            print(f"[DEBUG] Texte contenant conversion/taux (variations): '{text}'")
            # Extraire les variations LY/LW avec patterns plus robustes
            # Pattern LY plus flexible - capturer le % directement avant "VS LY"
            ly_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+ly',      # "+12% VS LY"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+ly',           # "+12% LY"
                r'vs\s+ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LY: +12%"
                r'ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LY +12%"
            ]
            
            for pattern in ly_patterns:
                ly_match = re.search(pattern, text_lower)
                if ly_match:
                    try:
                        cpfr_data['vs_ly_cr'] = float(ly_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] Conversion vs LY: {cpfr_data['vs_ly_cr']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LY: {e}")
                        continue
            
            # Pattern LW plus flexible - capturer le % directement avant "VS LW"
            lw_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+lw',      # "-14% VS LW"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+lw',           # "-14% LW"
                r'vs\s+lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LW: -14%"
                r'lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LW -14%"
            ]
            
            for pattern in lw_patterns:
                lw_match = re.search(pattern, text_lower)
                if lw_match:
                    try:
                        cpfr_data['vs_lw_cr'] = float(lw_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] Conversion vs LW: {cpfr_data['vs_lw_cr']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LW: {e}")
                        continue
        
        # BOOKINGS - traiter séparément la valeur et les variations
        # Pattern pour la valeur: "2 475"
        if text.strip().replace(' ', '').isdigit() and len(text.strip().replace(' ', '')) >= 3 and not any(word in text_lower for word in ['€', '%', 'k', 'm']):
            print(f"[DEBUG] Texte contenant nombre (valeur bookings): '{text}'")
            
            # Extraire la valeur principale - patterns pour capturer 2475 ou "2 475"
            booking_patterns = [
                r'(\d+(?:\s+\d+)*)',                            # "2 475" ou "2475"
                r'(\d{3,5})',                                   # nombre entre 3-5 chiffres
            ]
            
            for pattern in booking_patterns:
                booking_match = re.search(pattern, text.strip())
                if booking_match:
                    try:
                        value = booking_match.group(1).replace(' ', '')
                        if value.isdigit():
                            # Vérifier que c'est une valeur cohérente pour des bookings (100-10000)
                            amount = int(value)
                            if 100 <= amount <= 10000:
                                cpfr_data['nb_bookings'] = amount
                                print(f"[DEBUG] Bookings trouvés: {cpfr_data['nb_bookings']} (pattern: {pattern})")
                                break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing bookings: {e}")
                        continue
            
        # BOOKINGS - traiter les variations: "Nb of bookings +29% VS LY -18% VS LW"
        if 'booking' in text_lower or 'réservation' in text_lower:
            print(f"[DEBUG] Texte contenant booking/réservation (variations): '{text}'")
            # Extraire les variations LY/LW avec patterns plus robustes
            # Pattern LY plus flexible - capturer le % directement avant "VS LY"
            ly_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+ly',      # "+29% VS LY"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+ly',           # "+29% LY"
                r'vs\s+ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LY: +29%"
                r'ly[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LY +29%"
            ]
            
            for pattern in ly_patterns:
                ly_match = re.search(pattern, text_lower)
                if ly_match:
                    try:
                        cpfr_data['vs_ly_bookings'] = float(ly_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] Bookings vs LY: {cpfr_data['vs_ly_bookings']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LY: {e}")
                        continue
            
            # Pattern LW plus flexible - capturer le % directement avant "VS LW"
            lw_patterns = [
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+vs\s+lw',      # "-18% VS LW"
                r'([+-]?\d+(?:[,\.]\d+)*)\s*%\s+lw',           # "-18% LW"
                r'vs\s+lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',   # "VS LW: -18%"
                r'lw[:\s]*([+-]?\d+(?:[,\.]\d+)*)\s*%',        # "LW -18%"
            ]
            
            for pattern in lw_patterns:
                lw_match = re.search(pattern, text_lower)
                if lw_match:
                    try:
                        cpfr_data['vs_lw_bookings'] = float(lw_match.group(1).replace(',', '.')) / 100
                        print(f"[DEBUG] Bookings vs LW: {cpfr_data['vs_lw_bookings']} (pattern: {pattern})")
                        break
                    except Exception as e:
                        print(f"[DEBUG] Erreur parsing LW: {e}")
                        continue
        
        # Extraire tout nombre avec K, M, %, € pour analyse générale
        all_numbers = re.findall(r'(\d+(?:[,\.]\d+)*)\s*([km%€]?)', text_lower)
        if all_numbers:
            print(f"[DEBUG] Tous les nombres trouvés: {all_numbers}")
    
    print(f"[DEBUG] === RESULTAT FINAL CPFR ===")
    print(f"[DEBUG] Revenue final: {cpfr_data.get('revenue_b2c', 'NON DEFINI')}")
    print(f"[DEBUG] Résultat final CPFR: {cpfr_data}")
    return cpfr_data


def extract_cpfr_data_from_slide32(slide):
    """
    Extrait les données de tableau de la slide 32
    """
    table_data = parse_table(slide)
    
    cpfr_table_data = {
        'headers': table_data.get('headers', []),
        'rows': table_data.get('rows', []),
        'structured_data': {}
    }
    
    # Structurer les données du tableau
    if table_data.get('rows'):
        for row in table_data['rows']:
            if len(row) >= 2:
                key = row[0].strip().lower()
                value = row[1].strip()
                
                # Mapping des données du tableau
                if 'last minute' in key:
                    percentages = re.findall(r'(\d+[.,]?\d*)%', value)
                    if percentages:
                        cpfr_table_data['structured_data']['last_minute_pct'] = float(percentages[0].replace(',', '.')) / 100
                
                elif 'early booking' in key:
                    percentages = re.findall(r'(\d+[.,]?\d*)%', value)
                    if percentages:
                        cpfr_table_data['structured_data']['early_booking_pct'] = float(percentages[0].replace(',', '.')) / 100
                
                elif 'july' in key or 'juillet' in key:
                    percentages = re.findall(r'(\d+[.,]?\d*)%', value)
                    if percentages:
                        cpfr_table_data['structured_data']['month_july_pct'] = float(percentages[0].replace(',', '.')) / 100
                
                elif 'august' in key or 'août' in key:
                    percentages = re.findall(r'(\d+[.,]?\d*)%', value)
                    if percentages:
                        cpfr_table_data['structured_data']['month_august_pct'] = float(percentages[0].replace(',', '.')) / 100
                
                elif 'september' in key or 'septembre' in key:
                    percentages = re.findall(r'(\d+[.,]?\d*)%', value)
                    if percentages:
                        cpfr_table_data['structured_data']['month_sept_pct'] = float(percentages[0].replace(',', '.')) / 100
    
    return cpfr_table_data


def extract_cpfr_pptx(path, slide_start=31, slide_end=32):
    """
    Extrait spécifiquement les données CPFR des slides 31 et 32
    et les structure pour l'affichage groupé
    """
    try:
        prs = Presentation(path)
        slides = prs.slides
        
        # Validation des indices de slides
        if slide_start < 1 or slide_end < 1:
            raise ValueError("Les numéros de slides doivent être positifs")
        
        if slide_start > len(slides) or slide_end > len(slides):
            raise ValueError(f"Le fichier ne contient que {len(slides)} slides")
        
        if slide_start > slide_end:
            raise ValueError("La slide de début doit être inférieure à la slide de fin")
        
        # Extraction des slides
        slide31 = slides[slide_start - 1]
        slide32 = slides[slide_end - 1]
        
        # Extraction spécifique CPFR
        cpfr_data = extract_cpfr_data_from_slide31(slide31)
        table_data = extract_cpfr_data_from_slide32(slide32)
        
        # Fusion des données
        # Priorité aux données de la slide 31, puis complément avec la slide 32
        for key, value in table_data.get('structured_data', {}).items():
            if cpfr_data.get(key) is None:
                cpfr_data[key] = value
        
        # Structurer les données pour l'affichage groupé
        structured_preview = {
            'slide31_groups': [
                {
                    'label': 'Sessions',
                    'main_value': f"{cpfr_data.get('sessions', 0)//1000}K" if cpfr_data.get('sessions') and cpfr_data.get('sessions') >= 1000 else (f"{cpfr_data.get('sessions', 0):,}".replace(',', ' ') if cpfr_data.get('sessions') else 'Non trouvé'),
                    'unit': 'sessions',
                    'vs_ly': f"{cpfr_data.get('vs_ly_sessions', 0)*100:+.1f}%" if cpfr_data.get('vs_ly_sessions') else 'N/A',
                    'vs_lw': f"{cpfr_data.get('vs_lw_sessions', 0)*100:+.1f}%" if cpfr_data.get('vs_lw_sessions') else 'N/A',
                    'raw_texts': []  # Pour debug
                },
                {
                    'label': 'Web B2C Global revenue',
                    'main_value': f"{cpfr_data.get('revenue_b2c', 0)/1000000:.2f}M€".replace('.', ',') if cpfr_data.get('revenue_b2c') and cpfr_data.get('revenue_b2c') >= 1000000 else (f"{cpfr_data.get('revenue_b2c', 0):,.0f}€".replace(',', ' ') if cpfr_data.get('revenue_b2c') else 'Non trouvé'),
                    'unit': 'revenue',
                    'vs_ly': f"{cpfr_data.get('vs_ly_revenue', 0)*100:+.1f}%" if cpfr_data.get('vs_ly_revenue') else 'N/A',
                    'vs_lw': f"{cpfr_data.get('vs_lw_revenue', 0)*100:+.1f}%" if cpfr_data.get('vs_lw_revenue') else 'N/A',
                    'raw_texts': []
                },
                {
                    'label': 'Average Basket Value',
                    'main_value': f"{cpfr_data.get('average_basket_value', 0):.0f}€" if cpfr_data.get('average_basket_value') else 'Non trouvé',
                    'unit': 'basket/panier',
                    'vs_ly': f"{cpfr_data.get('vs_ly_abv', 0)*100:+.1f}%" if cpfr_data.get('vs_ly_abv') else 'N/A',
                    'vs_lw': f"{cpfr_data.get('vs_lw_abv', 0)*100:+.1f}%" if cpfr_data.get('vs_lw_abv') else 'N/A',
                    'raw_texts': []
                },
                {
                    'label': 'Conversion Rate',
                    'main_value': f"{cpfr_data.get('conversion_rate', 0)*100:.2f}%" if cpfr_data.get('conversion_rate') else 'Non trouvé',
                    'unit': 'conversion/taux',
                    'vs_ly': f"{cpfr_data.get('vs_ly_cr', 0)*100:+.2f}%" if cpfr_data.get('vs_ly_cr') else 'N/A',
                    'vs_lw': f"{cpfr_data.get('vs_lw_cr', 0)*100:+.2f}%" if cpfr_data.get('vs_lw_cr') else 'N/A',
                    'raw_texts': []
                },
                {
                    'label': 'Bookings',
                    'main_value': f"{cpfr_data.get('nb_bookings', 0):,}".replace(',', ' ') if cpfr_data.get('nb_bookings') else 'Non trouvé',
                    'unit': 'bookings/réservations',
                    'vs_ly': f"{cpfr_data.get('vs_ly_bookings', 0)*100:+.1f}%" if cpfr_data.get('vs_ly_bookings') else 'N/A',
                    'vs_lw': f"{cpfr_data.get('vs_lw_bookings', 0)*100:+.1f}%" if cpfr_data.get('vs_lw_bookings') else 'N/A',
                    'raw_texts': []
                }
            ]
        }
        
        # Ajouter les textes bruts trouvés pour debug - améliorer l'association
        all_texts = parse_slide_text(slide31)
        for text in all_texts:
            text_lower = text.lower()
            
            # Sessions - chercher patterns plus larges
            if any(word in text_lower for word in ['session', 'nb of sessions', 'sessions', '342k']):
                structured_preview['slide31_groups'][0]['raw_texts'].append(text)
            
            # Revenue - chercher patterns plus larges pour associer tous les textes liés
            if any(word in text_lower for word in ['revenue', 'web b2c', 'global revenue', '2,27m', '€']) or \
               any(word in text for word in ['2,27M€', '€2,27M']) or \
               ('€' in text and 'm' in text_lower):
                structured_preview['slide31_groups'][1]['raw_texts'].append(text)
            
            # Basket - chercher patterns plus larges
            if any(word in text_lower for word in ['basket', 'panier', 'average basket', 'panier moyen', '917€', '917', 'basket value']):
                structured_preview['slide31_groups'][2]['raw_texts'].append(text)
            
            # Conversion - chercher patterns plus larges
            if any(word in text_lower for word in ['conversion', 'taux', 'conversion rate', 'taux de conversion', '0,53%', '0.53%', 'rate']):
                structured_preview['slide31_groups'][3]['raw_texts'].append(text)
            
            # Bookings - chercher patterns plus larges
            if any(word in text_lower for word in ['booking', 'réservation', 'bookings', 'réservations', 'nb of bookings', '2475', 'do 2475']):
                structured_preview['slide31_groups'][4]['raw_texts'].append(text)
            
            # Capturer aussi les textes avec variations LY/LW
            if any(word in text_lower for word in ['vs ly', 'vs lw', '+%', '-%', 'ly', 'lw']) or \
               any(pattern in text for pattern in ['+11%', '-12%', '+6%', '-4%', '-15%', '+8%', '+12%', '-14%', '+29%', '-18%']):
                # Associer aux groupes selon le contexte
                if 'session' in text_lower:
                    structured_preview['slide31_groups'][0]['raw_texts'].append(text)
                elif 'revenue' in text_lower or 'b2c' in text_lower or 'global' in text_lower:
                    structured_preview['slide31_groups'][1]['raw_texts'].append(text)
                elif 'basket' in text_lower or 'panier' in text_lower or 'average' in text_lower:
                    structured_preview['slide31_groups'][2]['raw_texts'].append(text)
                elif 'conversion' in text_lower or 'taux' in text_lower or 'rate' in text_lower:
                    structured_preview['slide31_groups'][3]['raw_texts'].append(text)
                elif 'booking' in text_lower or 'réservation' in text_lower or 'nb of' in text_lower:
                    structured_preview['slide31_groups'][4]['raw_texts'].append(text)
        
        # Supprimer les doublons dans chaque groupe
        for group in structured_preview['slide31_groups']:
            group['raw_texts'] = list(dict.fromkeys(group['raw_texts']))  # Supprime doublons en gardant l'ordre
        
        return cpfr_data, table_data, structured_preview
        
    except Exception as e:
        raise Exception(f"Erreur lors de l'extraction CPFR du PowerPoint: {str(e)}")


def extract_pptx(path, slide_start, slide_end):
    """
    Extrait les données d'un fichier PowerPoint
    
    Args:
        path: Chemin vers le fichier .pptx
        slide_start: Numéro de la slide de début (1-indexed)
        slide_end: Numéro de la slide de fin (1-indexed)
    
    Returns:
        tuple: (kpis, table_data)
    """
    try:
        prs = Presentation(path)
        slides = prs.slides
        
        # Validation des indices de slides
        if slide_start < 1 or slide_end < 1:
            raise ValueError("Les numéros de slides doivent être positifs")
        
        if slide_start > len(slides) or slide_end > len(slides):
            raise ValueError(f"Le fichier ne contient que {len(slides)} slides")
        
        if slide_start > slide_end:
            raise ValueError("La slide de début doit être inférieure à la slide de fin")
        
        # Extraction des slides
        kpi_slide = slides[slide_start - 1]
        table_slide = slides[slide_end - 1]
        
        # Extraction des textes
        raw_texts = parse_slide_text(kpi_slide)
        kpis = extract_kpis_from_text(raw_texts)
        
        # Extraction du tableau
        table_data = parse_table(table_slide)
        
        return kpis, table_data
        
    except Exception as e:
        raise Exception(f"Erreur lors de l'extraction du PowerPoint: {str(e)}")


def get_slide_info(path):
    """Retourne des informations sur le fichier PowerPoint"""
    try:
        prs = Presentation(path)
        
        # Conversion des dates en chaînes pour la sérialisation JSON
        created_date = prs.core_properties.created
        modified_date = prs.core_properties.modified
        
        # Conversion en chaîne ISO si la date existe
        created_str = created_date.isoformat() if created_date else None
        modified_str = modified_date.isoformat() if modified_date else None
        
        return {
            "total_slides": len(prs.slides),
            "title": prs.core_properties.title or "Sans titre",
            "author": prs.core_properties.author or "Auteur inconnu",
            "created": created_str,
            "modified": modified_str
        }
    except Exception as e:
        raise Exception(f"Erreur lors de la lecture des informations du fichier: {str(e)}")
